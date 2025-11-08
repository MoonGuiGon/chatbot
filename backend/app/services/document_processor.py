"""
문서 처리 파이프라인
PDF, PPT, EXCEL, WORD 파일 파싱 및 처리
"""
import os
import uuid
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path

# Document parsers
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
from PIL import Image

from app.services.llm_service import get_embedding_llm, get_vision_llm
from app.config import config


@dataclass
class DocumentChunk:
    """문서 청크"""
    chunk_index: int
    content: str
    chunk_type: str  # text, table, image
    metadata: Dict[str, Any]
    image_path: Optional[str] = None
    embedding: Optional[List[float]] = None


@dataclass
class ParsedDocument:
    """파싱된 문서"""
    document_id: str
    file_name: str
    file_type: str
    chunks: List[DocumentChunk]
    metadata: Dict[str, Any]


class DocumentParser:
    """문서 파서 - 파일 타입별 파싱"""

    @staticmethod
    def parse(file_path: str) -> ParsedDocument:
        """파일 파싱"""
        file_ext = Path(file_path).suffix.lower()
        file_name = Path(file_path).name

        # 문서 ID 생성
        document_id = f"doc_{uuid.uuid4().hex[:12]}"

        if file_ext == ".pdf":
            chunks, metadata = DocumentParser._parse_pdf(file_path)
        elif file_ext in [".ppt", ".pptx"]:
            chunks, metadata = DocumentParser._parse_ppt(file_path)
        elif file_ext in [".doc", ".docx"]:
            chunks, metadata = DocumentParser._parse_docx(file_path)
        elif file_ext in [".xls", ".xlsx"]:
            chunks, metadata = DocumentParser._parse_excel(file_path)
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {file_ext}")

        return ParsedDocument(
            document_id=document_id,
            file_name=file_name,
            file_type=file_ext[1:],
            chunks=chunks,
            metadata=metadata
        )

    @staticmethod
    def _parse_pdf(file_path: str) -> Tuple[List[DocumentChunk], Dict[str, Any]]:
        """PDF 파싱"""
        reader = PdfReader(file_path)
        chunks = []

        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()

            if text.strip():
                chunk = DocumentChunk(
                    chunk_index=len(chunks),
                    content=text,
                    chunk_type="text",
                    metadata={
                        "page_number": page_num,
                        "source": "pdf"
                    }
                )
                chunks.append(chunk)

        metadata = {
            "total_pages": len(reader.pages),
            "file_type": "pdf"
        }

        return chunks, metadata

    @staticmethod
    def _parse_ppt(file_path: str) -> Tuple[List[DocumentChunk], Dict[str, Any]]:
        """PPT 파싱"""
        prs = Presentation(file_path)
        chunks = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts = []

            # 텍스트 추출
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)

            text = "\n".join(text_parts)

            if text.strip():
                chunk = DocumentChunk(
                    chunk_index=len(chunks),
                    content=text,
                    chunk_type="text",
                    metadata={
                        "slide_number": slide_num,
                        "source": "ppt"
                    }
                )
                chunks.append(chunk)

        metadata = {
            "total_slides": len(prs.slides),
            "file_type": "ppt"
        }

        return chunks, metadata

    @staticmethod
    def _parse_docx(file_path: str) -> Tuple[List[DocumentChunk], Dict[str, Any]]:
        """DOCX 파싱"""
        doc = Document(file_path)
        chunks = []

        current_chunk = []
        current_length = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                current_chunk.append(text)
                current_length += len(text)

                # 청크 크기 제한
                if current_length >= config.chunk_size:
                    chunk = DocumentChunk(
                        chunk_index=len(chunks),
                        content="\n".join(current_chunk),
                        chunk_type="text",
                        metadata={
                            "source": "docx"
                        }
                    )
                    chunks.append(chunk)
                    current_chunk = []
                    current_length = 0

        # 마지막 청크
        if current_chunk:
            chunk = DocumentChunk(
                chunk_index=len(chunks),
                content="\n".join(current_chunk),
                chunk_type="text",
                metadata={
                    "source": "docx"
                }
            )
            chunks.append(chunk)

        # 표 추출
        for table_num, table in enumerate(doc.tables, 1):
            table_text = DocumentParser._extract_table_text(table)
            chunk = DocumentChunk(
                chunk_index=len(chunks),
                content=table_text,
                chunk_type="table",
                metadata={
                    "table_number": table_num,
                    "source": "docx"
                }
            )
            chunks.append(chunk)

        metadata = {
            "total_paragraphs": len(doc.paragraphs),
            "total_tables": len(doc.tables),
            "file_type": "docx"
        }

        return chunks, metadata

    @staticmethod
    def _parse_excel(file_path: str) -> Tuple[List[DocumentChunk], Dict[str, Any]]:
        """Excel 파싱"""
        workbook = load_workbook(file_path)
        chunks = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            # 시트 데이터를 텍스트로 변환
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell else "" for cell in row])
                if row_text.strip():
                    rows.append(row_text)

            if rows:
                chunk = DocumentChunk(
                    chunk_index=len(chunks),
                    content="\n".join(rows),
                    chunk_type="table",
                    metadata={
                        "sheet_name": sheet_name,
                        "source": "excel"
                    }
                )
                chunks.append(chunk)

        metadata = {
            "total_sheets": len(workbook.sheetnames),
            "file_type": "excel"
        }

        return chunks, metadata

    @staticmethod
    def _extract_table_text(table) -> str:
        """표를 텍스트로 변환"""
        rows = []
        for row in table.rows:
            row_text = "\t".join([cell.text for cell in row.cells])
            rows.append(row_text)
        return "\n".join(rows)


class DocumentChunker:
    """문서 청킹 - Semantic Chunking"""

    @staticmethod
    def chunk(parsed_doc: ParsedDocument) -> ParsedDocument:
        """
        문서를 의미 단위로 청킹
        현재는 간단한 크기 기반 청킹
        """
        # 이미 청크된 상태이므로 추가 처리 없이 반환
        # 실제로는 더 정교한 Semantic Chunking 구현 필요

        # 청크 메타데이터 보강
        for i, chunk in enumerate(parsed_doc.chunks):
            chunk.metadata.update({
                "file_name": parsed_doc.file_name,
                "document_id": parsed_doc.document_id,
                "chunk_index": i,
                "total_chunks": len(parsed_doc.chunks)
            })

        return parsed_doc


class DocumentEmbedder:
    """문서 임베딩"""

    @staticmethod
    def embed(parsed_doc: ParsedDocument) -> ParsedDocument:
        """문서 청크 임베딩"""
        embedding_llm = get_embedding_llm()

        # 텍스트 수집
        texts = [chunk.content for chunk in parsed_doc.chunks]

        # 임베딩 생성
        embeddings = embedding_llm.embed_documents(texts)

        # 청크에 임베딩 추가
        for chunk, embedding in zip(parsed_doc.chunks, embeddings):
            chunk.embedding = embedding

        return parsed_doc


class DocumentProcessor:
    """
    문서 처리 파이프라인
    파싱 → 청킹 → 임베딩
    """

    def __init__(self):
        self.parser = DocumentParser()
        self.chunker = DocumentChunker()
        self.embedder = DocumentEmbedder()

    def process(self, file_path: str) -> ParsedDocument:
        """
        문서 처리 전체 파이프라인

        Args:
            file_path: 문서 파일 경로

        Returns:
            파싱 및 임베딩된 문서
        """
        # 1. 파싱
        parsed_doc = self.parser.parse(file_path)

        # 2. 청킹
        parsed_doc = self.chunker.chunk(parsed_doc)

        # 3. 임베딩
        parsed_doc = self.embedder.embed(parsed_doc)

        return parsed_doc

    def process_for_review(self, file_path: str) -> Dict[str, Any]:
        """
        검수용 문서 처리 (임베딩 제외)
        사용자가 검수 후 승인하면 임베딩 진행

        Args:
            file_path: 문서 파일 경로

        Returns:
            검수용 데이터
        """
        # 1. 파싱
        parsed_doc = self.parser.parse(file_path)

        # 2. 청킹
        parsed_doc = self.chunker.chunk(parsed_doc)

        # 3. 검수용 데이터 구성
        review_data = {
            "document_id": parsed_doc.document_id,
            "file_name": parsed_doc.file_name,
            "file_type": parsed_doc.file_type,
            "total_chunks": len(parsed_doc.chunks),
            "metadata": parsed_doc.metadata,
            "chunks": [
                {
                    "chunk_index": chunk.chunk_index,
                    "content": chunk.content,
                    "chunk_type": chunk.chunk_type,
                    "metadata": chunk.metadata,
                    "approved": True  # 기본값: 승인
                }
                for chunk in parsed_doc.chunks
            ]
        }

        return review_data

    def finalize_document(
        self,
        document_id: str,
        approved_chunks: List[Dict[str, Any]]
    ) -> ParsedDocument:
        """
        검수 완료된 문서 최종 처리 (임베딩)

        Args:
            document_id: 문서 ID
            approved_chunks: 승인된 청크 목록

        Returns:
            임베딩된 문서
        """
        # 승인된 청크만 필터링
        filtered_chunks = [
            chunk for chunk in approved_chunks
            if chunk.get("approved", False)
        ]

        # DocumentChunk 객체로 변환
        chunks = []
        for chunk_data in filtered_chunks:
            chunk = DocumentChunk(
                chunk_index=chunk_data["chunk_index"],
                content=chunk_data["content"],
                chunk_type=chunk_data["chunk_type"],
                metadata=chunk_data["metadata"]
            )
            chunks.append(chunk)

        # ParsedDocument 생성
        parsed_doc = ParsedDocument(
            document_id=document_id,
            file_name=filtered_chunks[0]["metadata"]["file_name"],
            file_type=filtered_chunks[0]["metadata"].get("source", "unknown"),
            chunks=chunks,
            metadata={}
        )

        # 임베딩
        parsed_doc = self.embedder.embed(parsed_doc)

        return parsed_doc
